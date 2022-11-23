from pptx import Presentation
from pptx.util import Cm
from datetime import date

def project_cost(df):
    d = df.drop_duplicates(subset='ProjectID')
    cost = sum(d['Cost'])
    cost_lastmonth = sum(d['Cost_lastmonth'])
    return cost, cost_lastmonth

def make_ppt(df):
    prs = Presentation()
    sld0 = prs.slides.add_slide(prs.slide_layouts[6])

    nrow, ncol = df.shape
    ppt_nrow = nrow + 2 # 列名と合計コスト挿入用に2行追加

    # shapeオブジェクト(表)を追加
    table_shape = sld0.shapes.add_table(ppt_nrow, ncol, Cm(1), Cm(1), Cm(24), Cm(15))
    table = table_shape.table
    tbl = table._graphic_frame.element.graphic.graphicData.tbl
    style_id = '{5940675A-B579-460E-94D1-54222C63F5DA}' #UUIDはgithub参照 https://github.com/scanny/python-pptx/issues/27
    tbl[0][-1].text = style_id

    # 行見出しのテキスト設定
    colnames = ['Project ID', 'インスタンスタイプ', 'Price', '今月のコスト', '先月のコスト']
    for i in range(len(colnames)):
        cell = table.cell(0, i)	 
        cell.text = colnames[i]

    # セルに値を設定する
    for i in range(nrow):
        value = list(df.iloc[i,])
        # 1つ上の行のProjectIDと、挿入する行のProjectIDが一致するなら、ProjectIDとProjectのコストは挿入しない
        if i == 0:
            former_value = ['', '', '', '']
        else:
            former_value = list(df.iloc[i-1,])
        if value[0] == former_value[0]:
            value[0] = ''
            value[3] = ''
            value[4] = ''

        for k in range(len(value)):
            cell = table.cell(i+1, k)
            cell.text = str(value[k])


    cost, cost_lastmonth = project_cost(df)
    cell = table.cell(ppt_nrow-1, 3)
    cell.text = 'Total: ' + str(cost) + '前月比: ' + str(round(cost/cost_lastmonth, 2)) + '倍'
    cell = table.cell(ppt_nrow-1, 4)
    cell.text = 'Total: ' + str(cost_lastmonth)

    # filename
    now = now = date.today()
    filename = now.strftime('%Y%m%d') + '_powerpoint_name.pptx'
    prs.save(filename)

    return('The function has completed')

